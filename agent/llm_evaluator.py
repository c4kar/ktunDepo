"""
ktunDepo Agent — LLM Evaluator
OpenRouter API üzerinden materyal kalite değerlendirmesi yapar.

Sadece heuristik filtrenin "gri bölge" kararlarında (NEEDS_LLM) devreye girer.
Token tasarrufu için ilk 3 sayfadan ~3000 kelime sınırıyla metin alır.
"""

import os
import json
import time
import httpx
from pathlib import Path
from dataclasses import dataclass
from typing import Optional, Tuple

from agent.logging_config import get_logger

logger = get_logger("llm_evaluator")

# OpenRouter API endpoint
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"

# Model: hızlı ve ucuz, Türkçe anlayışı iyi
DEFAULT_MODEL = "anthropic/claude-3-haiku-20240307"
FALLBACK_MODEL = "openai/gpt-4o-mini"

# Metin limiti (token tasarrufu)
MAX_TEXT_CHARS = 8000  # ~2000 token kadar metin


@dataclass
class LLMEvalResult:
    """LLM değerlendirme sonucu."""

    decision: str  # ACCEPT | REJECT | SIMILAR_EXISTS
    confidence: float  # 0.0 - 1.0
    reason: str
    detected_course: Optional[str] = None
    material_type: Optional[str] = None
    suggested_path: Optional[str] = None
    tokens_used: int = 0
    model_used: str = ""
    error: Optional[str] = None


class LLMEvaluator:
    """
    OpenRouter API ile materyal kalite değerlendirmesi.

    Kullanım akışı:
    1. Dosyadan metin çıkar (PDF/PPTX/DOCX)
    2. Prompt şablonunu yükle
    3. OpenRouter'a gönder
    4. JSON yanıtı parse et
    """

    def __init__(
        self,
        api_key: Optional[str] = None,
        model: str = DEFAULT_MODEL,
        prompts_path: str = "agent/prompts",
        timeout: int = 30,
    ):
        self.api_key = api_key or os.getenv("OPENROUTER_API_KEY", "")
        self.model = model
        self.prompts_path = Path(prompts_path)
        self.timeout = timeout
        self._system_prompt: Optional[str] = None

    def _load_system_prompt(self) -> str:
        """quality_evaluation.md prompt şablonunu yükle."""
        if self._system_prompt:
            return self._system_prompt

        prompt_file = self.prompts_path / "quality_evaluation.md"
        if prompt_file.exists():
            self._system_prompt = prompt_file.read_text(encoding="utf-8")
        else:
            # Hardcoded fallback
            self._system_prompt = (
                "Sen bir üniversite ders materyali kalite değerlendiricisisin.\n"
                "Materyalin depoya eklenip eklenmeyeceğine karar ver.\n\n"
                "Yanıt formatı (JSON):\n"
                '{"decision": "ACCEPT|REJECT", "confidence": 0.0-1.0, '
                '"reason": "kısa açıklama", "detected_course": null, '
                '"material_type": "Ders Notu|Sınav Sorusu|Sunum|Özet|Diğer", '
                '"suggested_path": null}'
            )
        return self._system_prompt

    def _extract_text(self, file_path: str) -> str:
        """Dosyadan metin çıkar (PDF, PPTX, DOCX)."""
        path = Path(file_path)
        ext = path.suffix.lower()
        text = ""

        if ext == ".pdf":
            text = self._extract_pdf_text(file_path)
        elif ext == ".pptx":
            text = self._extract_pptx_text(file_path)
        elif ext == ".docx":
            text = self._extract_docx_text(file_path)
        else:
            # Video veya desteklenmeyen format — sadece dosya adı gönder
            text = f"[Dosya adı: {path.name}] [Format: {ext}]"

        # Karakter limiti uygula
        return text[:MAX_TEXT_CHARS]

    def _extract_pdf_text(self, file_path: str) -> str:
        """PDF'den metin çıkar."""
        text = ""
        # pypdf2 dene
        try:
            from pypdf import PdfReader  # type: ignore

            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages[:5]:
                    text += (page.extract_text() or "") + "\n"
            if len(text.strip()) > 100:
                return text
        except Exception:
            pass

        # pdfplumber fallback
        try:
            import pdfplumber  # type: ignore

            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:5]:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"
        except Exception:
            pass

        return text

    def _extract_pptx_text(self, file_path: str) -> str:
        """PPTX'ten metin çıkar."""
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(file_path)
            text_parts = []
            for i, slide in enumerate(prs.slides[:10]):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                if i >= 9:
                    break
            return "\n".join(text_parts)
        except Exception:
            return ""

    def _extract_docx_text(self, file_path: str) -> str:
        """DOCX'ten metin çıkar."""
        try:
            from docx import Document  # type: ignore

            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs[:100] if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception:
            return ""

    def _call_api(
        self, system_prompt: str, user_message: str, model: str
    ) -> Tuple[Optional[dict], int]:
        """
        OpenRouter API'ye senkron çağrı yap.

        Returns:
            (parsed_response_dict, tokens_used)
        """
        if not self.api_key:
            raise ValueError("OPENROUTER_API_KEY ayarlanmamış")

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://github.com/c4kar/ktunDepo",
            "X-Title": "ktunDepo Quality Evaluator",
        }

        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message},
            ],
            "temperature": 0.1,  # Kararlı, düşük yaratıcılık
            "max_tokens": 400,
            "response_format": {"type": "json_object"},
        }

        with httpx.Client(timeout=self.timeout) as client:
            response = client.post(OPENROUTER_API_URL, headers=headers, json=payload)
            response.raise_for_status()
            data = response.json()

        tokens_used = data.get("usage", {}).get("total_tokens", 0)
        content = data["choices"][0]["message"]["content"]
        parsed = json.loads(content)
        return parsed, tokens_used

    def evaluate(
        self,
        file_path: str,
        course_hint: Optional[str] = None,
        material_type_hint: Optional[str] = None,
        heuristic_score: int = 50,
        duplicate_similar: bool = False,
    ) -> LLMEvalResult:
        """
        Materyali LLM ile değerlendir.

        Args:
            file_path: Değerlendirilecek dosya
            course_hint: Metadata'dan gelen ders adı (opsiyonel)
            material_type_hint: Metadata'dan gelen materyal türü (opsiyonel)
            heuristic_score: Heuristik kalite skoru (bağlam için)
            duplicate_similar: Benzer materyal tespit edildi mi (bağlam için)

        Returns:
            LLMEvalResult
        """
        if not self.api_key:
            logger.warning("OPENROUTER_API_KEY yok, LLM değerlendirmesi atlandı")
            # API key yoksa skora göre karar ver
            decision = "ACCEPT" if heuristic_score >= 45 else "REJECT"
            return LLMEvalResult(
                decision=decision,
                confidence=0.5,
                reason="LLM API anahtarı eksik, skor bazlı karar verildi",
                error="API key missing",
            )

        file_path_obj = Path(file_path)
        filename = file_path_obj.name

        # Metin çıkar
        extracted_text = self._extract_text(file_path)

        # Kullanıcı mesajı oluştur
        context_parts = [f"**Dosya adı:** {filename}"]
        if course_hint:
            context_parts.append(f"**Belirtilen ders:** {course_hint}")
        if material_type_hint:
            context_parts.append(f"**Belirtilen materyal türü:** {material_type_hint}")
        context_parts.append(f"**Heuristik skor:** {heuristic_score}/100")
        if duplicate_similar:
            context_parts.append(
                "**Not:** Vektör veritabanında benzer bir materyal tespit edildi."
            )

        if extracted_text.strip():
            context_parts.append(
                f"\n**İçerik (ilk bölüm):**\n```\n{extracted_text}\n```"
            )
        else:
            context_parts.append(
                "\n**Not:** Dosyadan metin çıkarılamadı (görsel/taranmış PDF olabilir)."
            )

        user_message = "\n".join(context_parts)
        system_prompt = self._load_system_prompt()

        # API çağrısı (fallback model ile retry)
        for attempt, model in enumerate([self.model, FALLBACK_MODEL]):
            try:
                parsed, tokens = self._call_api(system_prompt, user_message, model)

                decision = parsed.get("decision", "REJECT").upper()
                # "ACCEPT" → kabul, diğer her şey → red
                if decision not in ("ACCEPT", "REJECT", "SIMILAR_EXISTS"):
                    decision = "REJECT"

                return LLMEvalResult(
                    decision=decision,
                    confidence=float(parsed.get("confidence", 0.5)),
                    reason=parsed.get("reason", "LLM kararı"),
                    detected_course=parsed.get("detected_course"),
                    material_type=parsed.get("material_type"),
                    suggested_path=parsed.get("suggested_path"),
                    tokens_used=tokens,
                    model_used=model,
                )

            except httpx.HTTPStatusError as e:
                logger.warning(
                    f"OpenRouter HTTP hatası ({model}): {e.response.status_code}"
                )
                if attempt == 0:
                    time.sleep(1)
                    continue
                # Her iki model de başarısız → güvenli karar
                return LLMEvalResult(
                    decision="REJECT" if heuristic_score < 45 else "ACCEPT",
                    confidence=0.4,
                    reason=f"LLM API hatası: HTTP {e.response.status_code}",
                    error=str(e),
                )

            except json.JSONDecodeError as e:
                logger.warning(f"LLM JSON parse hatası ({model}): {e}")
                if attempt == 0:
                    continue
                return LLMEvalResult(
                    decision="REJECT" if heuristic_score < 45 else "ACCEPT",
                    confidence=0.4,
                    reason="LLM yanıtı parse edilemedi",
                    error=str(e),
                )

            except Exception as e:
                logger.error(f"LLM değerlendirme hatası ({model}): {e}")
                if attempt == 0:
                    time.sleep(1)
                    continue
                return LLMEvalResult(
                    decision="REJECT" if heuristic_score < 45 else "ACCEPT",
                    confidence=0.3,
                    reason=f"LLM hatası: {str(e)[:100]}",
                    error=str(e),
                )

        # Buraya gelmemeli ama type checker için
        return LLMEvalResult(
            decision="REJECT",
            confidence=0.0,
            reason="Bilinmeyen hata",
            error="unreachable",
        )


# Singleton
_evaluator: Optional[LLMEvaluator] = None


def get_llm_evaluator(
    api_key: Optional[str] = None,
    model: str = DEFAULT_MODEL,
    prompts_path: str = "agent/prompts",
) -> LLMEvaluator:
    """Global LLMEvaluator instance döndür."""
    global _evaluator
    if _evaluator is None:
        _evaluator = LLMEvaluator(
            api_key=api_key,
            model=model,
            prompts_path=prompts_path,
        )
    return _evaluator
