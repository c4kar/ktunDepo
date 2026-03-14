"""
ktunDepo Intake Agent — File Scanner
Teknik tarama modülü. Sadece fiziksel olarak bozuk dosyaları eliyor.

ÖNEMLI: Bu modül içerik kalitesi hakkında karar VERMEZ.
Sayfa sayısı, metin yoğunluğu gibi metrikler red kriteri DEĞİLDİR.
"""

import os
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional, Dict, Any, List
from enum import Enum


class TechnicalError(Exception):
    """Teknik tarama hatası — dosya işlenemez."""

    pass


class UnsupportedFormatError(TechnicalError):
    """Desteklenmeyen dosya formatı."""

    pass


class EmptyFileError(TechnicalError):
    """Dosya boş veya çok küçük."""

    pass


class CorruptedFileError(TechnicalError):
    """Dosya bozuk, açılamıyor."""

    pass


@dataclass
class ScanResult:
    """Teknik tarama sonucu."""

    filename_original: str
    extension: str
    size_kb: float
    size_mb: float
    file_path: str

    # PDF'e özgü (varsa)
    page_count: Optional[int] = None
    has_text_layer: bool = False
    first_page_text: str = ""
    first_page_has_images: bool = False

    # PPTX/DOCX için
    slide_count: Optional[int] = None

    # Görüntü dosyaları
    is_image: bool = False

    # Video dosyaları
    is_video: bool = False

    # Metadata
    extra: Dict[str, Any] = field(default_factory=dict)


class FileScanner:
    """
    Teknik dosya tarama sınıfı.

    Sadece şunları kontrol eder:
    - Dosya var mı ve açılabiliyor mu?
    - Format destekleniyor mu?
    - Dosya boş mu (< 5KB)?

    İçerik kalitesi hakkında karar VERMEZ.
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".pptx",
        ".docx",
        ".jpg",
        ".jpeg",
        ".png",
        ".mp4",
        ".mkv",
    }

    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}
    VIDEO_EXTENSIONS = {".mp4", ".mkv"}

    MIN_FILE_SIZE_KB = 5  # 5KB altı dosyalar muhtemelen boş/bozuk

    def scan(self, file_path: str) -> ScanResult:
        """
        Dosyayı teknik olarak tara.

        Args:
            file_path: Taranacak dosya yolu

        Returns:
            ScanResult objesi

        Raises:
            TechnicalError: Dosya teknik olarak işlenemezse
        """
        path = Path(file_path)

        # Dosya var mı?
        if not path.exists():
            raise CorruptedFileError(f"Dosya bulunamadı: {file_path}")

        # Temel bilgiler
        stat = path.stat()
        size_bytes = stat.st_size
        size_kb = size_bytes / 1024
        size_mb = size_bytes / (1024 * 1024)
        extension = path.suffix.lower()

        # Uzantı kontrolü
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise UnsupportedFormatError(
                f"Desteklenmeyen format: {extension}. "
                f"Desteklenen: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        # Boyut kontrolü — sadece gerçekten boş dosyalar
        if size_kb < self.MIN_FILE_SIZE_KB:
            raise EmptyFileError(
                f"Dosya {size_kb:.1f}KB, minimum {self.MIN_FILE_SIZE_KB}KB gerekli"
            )

        # Temel result
        result = ScanResult(
            filename_original=path.name,
            extension=extension,
            size_kb=size_kb,
            size_mb=size_mb,
            file_path=str(path.absolute()),
            is_image=extension in self.IMAGE_EXTENSIONS,
            is_video=extension in self.VIDEO_EXTENSIONS,
        )

        # Format'a özgü tarama
        if extension == ".pdf":
            self._scan_pdf(file_path, result)
        elif extension == ".pptx":
            self._scan_pptx(file_path, result)
        elif extension == ".docx":
            self._scan_docx(file_path, result)
        elif extension in self.IMAGE_EXTENSIONS:
            self._scan_image(file_path, result)
        # Video için ekstra tarama yok

        return result

    def _scan_pdf(self, file_path: str, result: ScanResult) -> None:
        """PDF'i tara — sayfa sayısı ve metin katmanı bilgisi."""
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                # PDF kütüphanesi yok — LLM vision ile devam edebilir
                result.extra["pdf_library_missing"] = True
                return

        try:
            reader = PdfReader(file_path)
            result.page_count = len(reader.pages)

            # Metin katmanı var mı? İlk 3 sayfayı dene
            sample_text = ""
            for i, page in enumerate(reader.pages[:3]):
                try:
                    text = page.extract_text() or ""
                    sample_text += text
                except Exception:
                    pass

            result.has_text_layer = len(sample_text.strip()) > 50

            # İlk sayfa metnini kaydet
            if result.has_text_layer and reader.pages:
                try:
                    result.first_page_text = reader.pages[0].extract_text() or ""
                except Exception:
                    pass

            # İlk sayfada görsel var mı?
            if reader.pages:
                try:
                    first_page = reader.pages[0]
                    if hasattr(first_page, "images"):
                        result.first_page_has_images = len(first_page.images) > 0
                except Exception:
                    pass

        except Exception as e:
            raise CorruptedFileError(f"PDF açılamadı: {e}")

    def _scan_pptx(self, file_path: str, result: ScanResult) -> None:
        """PPTX'i tara."""
        try:
            from pptx import Presentation
        except ImportError:
            result.extra["pptx_library_missing"] = True
            result.has_text_layer = True  # Office formatları metin içerir
            return

        try:
            prs = Presentation(file_path)
            result.slide_count = len(prs.slides)
            result.page_count = result.slide_count
            result.has_text_layer = True

            # İlk slayttan metin çıkar
            if prs.slides:
                texts = []
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                result.first_page_text = "\n".join(texts[:5])

        except Exception as e:
            raise CorruptedFileError(f"PPTX açılamadı: {e}")

    def _scan_docx(self, file_path: str, result: ScanResult) -> None:
        """DOCX'i tara."""
        try:
            from docx import Document
        except ImportError:
            result.extra["docx_library_missing"] = True
            result.has_text_layer = True
            return

        try:
            doc = Document(file_path)
            result.has_text_layer = True

            # İlk paragrafları al
            paragraphs = [p.text for p in doc.paragraphs[:10] if p.text.strip()]
            result.first_page_text = "\n".join(paragraphs)

            # Yaklaşık sayfa sayısı (her 3000 karakter = 1 sayfa)
            total_text = " ".join(p.text for p in doc.paragraphs)
            result.page_count = max(1, len(total_text) // 3000)

        except Exception as e:
            raise CorruptedFileError(f"DOCX açılamadı: {e}")

    def _scan_image(self, file_path: str, result: ScanResult) -> None:
        """Görüntü dosyasını tara."""
        try:
            from PIL import Image
        except ImportError:
            result.extra["pil_library_missing"] = True
            return

        try:
            with Image.open(file_path) as img:
                result.extra["image_size"] = img.size
                result.extra["image_mode"] = img.mode
                result.page_count = 1
                result.has_text_layer = False  # Görüntü = OCR gerekli

        except Exception as e:
            raise CorruptedFileError(f"Görüntü açılamadı: {e}")

    def scan_directory(self, directory: str) -> List[tuple]:
        """
        Bir klasördeki tüm dosyaları tara.

        Args:
            directory: Taranacak klasör

        Returns:
            List of (file_path, ScanResult or Exception) tuples
        """
        results = []
        dir_path = Path(directory)

        if not dir_path.exists():
            return results

        for file_path in dir_path.iterdir():
            if file_path.is_file() and not file_path.name.startswith("."):
                try:
                    result = self.scan(str(file_path))
                    results.append((str(file_path), result))
                except TechnicalError as e:
                    results.append((str(file_path), e))

        return results
