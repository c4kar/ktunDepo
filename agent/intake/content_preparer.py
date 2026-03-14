"""
ktunDepo Intake Agent — Content Preparer
LLM'e gönderilecek içeriği hazırlar.

Metin katmanı varsa → metin modu
Metin katmanı yoksa → vision modu (PDF sayfasını görüntüye çevir)
Görüntü dosyası → direkt vision modu
Video dosyası → sadece metadata
"""

import base64
import io
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any
from enum import Enum

from agent.intake.file_scanner import ScanResult


class ContentMode(Enum):
    """İçerik hazırlama modu."""

    TEXT = "text"  # Metin katmanı var
    VISION = "vision"  # Görsel analiz gerekli
    METADATA_ONLY = "metadata"  # Video — sadece dosya adı/boyut


@dataclass
class PreparedContent:
    """LLM'e gönderilecek hazırlanmış içerik."""

    mode: ContentMode

    # TEXT modu için
    text_content: str = ""
    pages_sampled: int = 0

    # VISION modu için
    images_base64: List[str] = field(default_factory=list)
    media_type: str = "image/jpeg"

    # METADATA_ONLY modu için
    filename: str = ""
    size_mb: float = 0.0

    # Genel
    note: str = ""
    extra: Dict[str, Any] = field(default_factory=dict)


class ContentPreparer:
    """
    LLM'e gönderilecek içeriği hazırlayan sınıf.

    Dosya türüne göre uygun modu seçer:
    - PDF/PPTX/DOCX + metin varsa → text modu
    - PDF metin yoksa → vision modu (pdf2image ile)
    - JPG/PNG → vision modu (direkt)
    - MP4/MKV → metadata modu
    """

    MAX_TEXT_CHARS = 3000  # Metin için max karakter
    MAX_PAGES_TEXT = 2  # Metin için max sayfa
    MAX_PAGES_VISION = 2  # Vision için max sayfa
    VISION_DPI = 200  # PDF → görüntü dönüşüm çözünürlüğü
    VISION_QUALITY = 75  # JPEG kalitesi (0-100). Daha düşük = daha küçük dosya, 75 yeterli

    def prepare(self, scan_result: ScanResult) -> PreparedContent:
        """
        Tarama sonucuna göre içerik hazırla.

        Args:
            scan_result: FileScanner'dan gelen tarama sonucu

        Returns:
            PreparedContent objesi
        """
        # Video dosyası — sadece metadata
        if scan_result.is_video:
            return self._prepare_video(scan_result)

        # Görüntü dosyası — vision modu
        if scan_result.is_image:
            return self._prepare_image(scan_result)

        # PDF — metin varsa text, yoksa vision
        if scan_result.extension == ".pdf":
            if scan_result.has_text_layer:
                return self._prepare_pdf_text(scan_result)
            else:
                return self._prepare_pdf_vision(scan_result)

        # PPTX/DOCX — metin modu
        if scan_result.extension in (".pptx", ".docx"):
            return self._prepare_office_text(scan_result)

        # Fallback — metadata
        return PreparedContent(
            mode=ContentMode.METADATA_ONLY,
            filename=scan_result.filename_original,
            size_mb=scan_result.size_mb,
            note="Bilinmeyen format, sadece metadata",
        )

    def _prepare_video(self, scan: ScanResult) -> PreparedContent:
        """Video dosyası için metadata hazırla."""
        return PreparedContent(
            mode=ContentMode.METADATA_ONLY,
            filename=scan.filename_original,
            size_mb=scan.size_mb,
            note="Video dosyası — içerik analizi yapılamaz, sadece metadata ile değerlendir",
        )

    def _prepare_image(self, scan: ScanResult) -> PreparedContent:
        """Görüntü dosyasını base64'e çevir. Boyutu optimize et."""
        try:
            from PIL import Image
            
            # Orijinal dosyayı aç ve yeniden kodla (kaliteyi düşür)
            with Image.open(scan.file_path) as img:
                buffer = io.BytesIO()
                # JPEG formatına dönüştür ve kaliteyi düşür
                img.convert('RGB').save(buffer, format='JPEG', quality=self.VISION_QUALITY)
                b64 = base64.b64encode(buffer.getvalue()).decode()
            
            media_type = "image/jpeg"
            
            return PreparedContent(
                mode=ContentMode.VISION,
                images_base64=[b64],
                media_type=media_type,
                note="Görüntü dosyası, vision ile analiz ediliyor",
            )

        except Exception as e:
            return PreparedContent(
                mode=ContentMode.METADATA_ONLY,
                filename=scan.filename_original,
                size_mb=scan.size_mb,
                note=f"Görüntü okunamadı: {e}",
            )

    def _prepare_pdf_text(self, scan: ScanResult) -> PreparedContent:
        """Metin katmanı olan PDF için metin çıkar."""
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                # Fallback to first_page_text from scan
                return PreparedContent(
                    mode=ContentMode.TEXT,
                    text_content=scan.first_page_text[: self.MAX_TEXT_CHARS],
                    pages_sampled=1,
                    note="PDF kütüphanesi eksik, scan'den alınan metin",
                )

        try:
            reader = PdfReader(scan.file_path)
            text_parts = []

            pages_to_read = min(self.MAX_PAGES_TEXT, len(reader.pages))
            for i in range(pages_to_read):
                page_text = reader.pages[i].extract_text() or ""
                text_parts.append(f"\n--- SAYFA {i + 1} ---\n{page_text}")

            combined = "".join(text_parts)[: self.MAX_TEXT_CHARS]

            return PreparedContent(
                mode=ContentMode.TEXT,
                text_content=combined,
                pages_sampled=pages_to_read,
                note=f"İlk {pages_to_read} sayfanın metni",
            )

        except Exception as e:
            # Metin çıkarma başarısız — vision'a geç
            return self._prepare_pdf_vision(scan)

    def _prepare_pdf_vision(self, scan: ScanResult) -> PreparedContent:
        """Metin katmanı olmayan PDF için görüntü hazırla."""
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return PreparedContent(
                mode=ContentMode.METADATA_ONLY,
                filename=scan.filename_original,
                size_mb=scan.size_mb,
                note="pdf2image kütüphanesi eksik, vision modu kullanılamıyor",
                extra={"needs_pdf2image": True},
            )

        try:
            page_count = scan.page_count or 1
            pages_to_convert = min(self.MAX_PAGES_VISION, page_count)

            pages = convert_from_path(
                scan.file_path,
                first_page=1,
                last_page=pages_to_convert,
                dpi=self.VISION_DPI,
            )

            images_b64 = []
            for page_img in pages:
                buffer = io.BytesIO()
                page_img.save(buffer, format="JPEG", quality=self.VISION_QUALITY)
                b64 = base64.b64encode(buffer.getvalue()).decode()
                images_b64.append(b64)

            return PreparedContent(
                mode=ContentMode.VISION,
                images_base64=images_b64,
                media_type="image/jpeg",
                note=f"Metin katmanı yok, ilk {pages_to_convert} sayfa görsel olarak analiz ediliyor",
                extra={"total_pages": page_count},
            )

        except Exception as e:
            return PreparedContent(
                mode=ContentMode.METADATA_ONLY,
                filename=scan.filename_original,
                size_mb=scan.size_mb,
                note=f"PDF görüntüye çevrilemedi: {e}",
            )

    def _prepare_office_text(self, scan: ScanResult) -> PreparedContent:
        """PPTX/DOCX için metin hazırla."""
        # Scan'den alınan first_page_text'i kullan
        text = scan.first_page_text

        # Eğer yetersizse daha fazla çıkarmayı dene
        if len(text) < 200:
            text = self._extract_more_text(scan)

        return PreparedContent(
            mode=ContentMode.TEXT,
            text_content=text[: self.MAX_TEXT_CHARS],
            pages_sampled=min(scan.page_count or 1, self.MAX_PAGES_TEXT),
            note=f"{scan.extension.upper()} dosyası",
        )

    def _extract_more_text(self, scan: ScanResult) -> str:
        """Office dosyasından daha fazla metin çıkar."""
        if scan.extension == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(scan.file_path)
                texts = []
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                return "\n".join(texts)
            except Exception:
                pass

        elif scan.extension == ".docx":
            try:
                from docx import Document

                doc = Document(scan.file_path)
                paragraphs = [p.text for p in doc.paragraphs[:30] if p.text.strip()]
                return "\n".join(paragraphs)
            except Exception:
                pass

        return scan.first_page_text
